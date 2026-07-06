import os
import base64

SUPPORTED_EXTENSIONS = {'pdf', 'txt', 'md', 'markdown', 'png', 'jpg', 'jpeg', 'webp', 'gif', 'docx', 'pptx', 'csv'}

IMAGE_MIME = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'webp': 'image/webp',
    'gif': 'image/gif',
}


def extract(file_path: str) -> dict:
    ext = file_path.rsplit('.', 1)[-1].lower()

    if ext == 'pdf':
        return _extract_pdf(file_path)
    elif ext in IMAGE_MIME:
        return _extract_image(file_path, ext)
    elif ext == 'docx':
        return _extract_docx(file_path)
    elif ext == 'pptx':
        return _extract_pptx(file_path)
    elif ext == 'csv':
        return _extract_csv(file_path)
    elif ext in ('txt', 'md', 'markdown'):
        return _extract_text(file_path)
    else:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'unsupported', 'error': f'Unsupported: {ext}'}


def _extract_pdf(path):
    try:
        import fitz
        doc = fitz.open(path)
        text = '\n'.join(page.get_text() for page in doc)
        doc.close()

        if len(text.strip()) > 100:
            return {'text': text, 'image_b64': None, 'mime_type': None, 'method': 'pdf_text', 'error': None}

        # Scanned PDF - use vision on first page
        doc = fitz.open(path)
        page = doc[0]
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes('png')
        doc.close()
        b64 = base64.b64encode(img_bytes).decode()
        return {'text': '', 'image_b64': b64, 'mime_type': 'image/png', 'method': 'pdf_vision', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'pdf_error', 'error': str(e)}


def _extract_image(path, ext):
    try:
        with open(path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode()
        return {'text': '', 'image_b64': b64, 'mime_type': IMAGE_MIME[ext], 'method': 'image_vision', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'image_error', 'error': str(e)}


def _extract_docx(path):
    """
    Extract text from a .docx file in document order, INCLUDING tables.

    Word RFPs and contracts put critical data (property overview, SLAs,
    staffing, evaluation criteria, pricing worksheets) in tables. Reading
    only doc.paragraphs silently drops all of that content and causes the
    AI to hallucinate that sections are missing when they are actually
    present in tabular form.
    """
    try:
        from docx import Document
        from docx.document import Document as _DocBody
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph

        def _iter_block_items(parent):
            if isinstance(parent, _DocBody):
                parent_elm = parent.element.body
            elif isinstance(parent, _Cell):
                parent_elm = parent._tc
            else:
                return
            for child in parent_elm.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)

        def _render_table(table):
            rows_out = []
            for row in table.rows:
                cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                if any(cells):
                    rows_out.append(' | '.join(cells))
            if not rows_out:
                return ''
            return '\n'.join(['[TABLE]', *rows_out, '[/TABLE]'])

        doc = Document(path)
        parts = []
        for block in _iter_block_items(doc):
            if isinstance(block, Paragraph):
                txt = block.text.strip()
                if txt:
                    parts.append(txt)
            elif isinstance(block, Table):
                rendered = _render_table(block)
                if rendered:
                    parts.append(rendered)

        text = '\n'.join(parts)
        return {'text': text, 'image_b64': None, 'mime_type': None, 'method': 'docx', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'docx_error', 'error': str(e)}


def _extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(run.text for run in para.runs).strip()
                        if txt:
                            lines.append(txt)
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[Speaker notes] {notes}")
            if len(lines) > 1:
                parts.append("\n".join(lines))
        text = "\n\n".join(parts)
        return {'text': text, 'image_b64': None, 'mime_type': None, 'method': 'pptx', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'pptx_error', 'error': str(e)}


def _extract_csv(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        return {'text': text, 'image_b64': None, 'mime_type': None, 'method': 'csv', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'csv_error', 'error': str(e)}


def _extract_text(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        return {'text': text, 'image_b64': None, 'mime_type': None, 'method': 'plaintext', 'error': None}
    except Exception as e:
        return {'text': '', 'image_b64': None, 'mime_type': None, 'method': 'text_error', 'error': str(e)}
